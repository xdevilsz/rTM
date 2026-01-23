import argparse
import asyncio
import hashlib
import hmac
import json
import os
import threading
import time
import urllib.parse
import urllib.request
from collections import deque
from datetime import datetime, timedelta, timezone
from http.server import SimpleHTTPRequestHandler
try:
    # Python 3.7+: handles concurrent requests (critical when /api/analysis is slow)
    from http.server import ThreadingHTTPServer as HTTPServer
except Exception:
    # Fallback (single-threaded)
    from http.server import HTTPServer
from pathlib import Path
from typing import Dict, List, Optional
from dotenv import load_dotenv

# Import our new modules
from analyzer import TradeAnalyzer
from bybit_client import BybitAPIClient, BybitWebSocketClient
from trading_service import TradingService

# Load .env file from current directory or script directory
load_dotenv()
load_dotenv(Path(__file__).resolve().parent / ".env")  # Also check script directory

ROOT = Path(__file__).resolve().parent
DASHBOARD_DIR = ROOT / "dashboard"
ADMIN_CONFIG_PATH = ROOT / "admin_config.json"
REPORTS_DIR = ROOT / "reports"
SERVER_START_TS = time.time()

# Global state for realtime mode
BROKER_ID = "To000948"
_realtime_data = {
    "fills": deque(maxlen=10000),  # Keep last 10k fills
    "non_trade_events": deque(maxlen=5000),
    "positions": {},
    "orders": deque(maxlen=1000),
    "last_update": time.time(),
    "running": False,
}
_realtime_lock = threading.Lock()


def _env(key: str, default: str = "") -> str:
    return os.getenv(key, default).strip()


def _to_float(value):
    try:
        if value is None or value == "":
            return None
        return float(value)
    except Exception:
        return None


def _safe_int(value):
    try:
        return int(value)
    except Exception:
        return 0


def _normalize_ts(value) -> float:
    """Normalize timestamps to seconds."""
    try:
        ts = float(value)
    except Exception:
        return 0.0
    if ts > 1e12:
        ts = ts / 1000.0
    return ts


def _build_metrics_from_fills(fills: list[dict]) -> dict:
    trade_fills = [f for f in fills if not f.get("exec_type") or f.get("exec_type") == "trade"]
    non_trade_events = [f for f in fills if f.get("exec_type") and f.get("exec_type") != "trade"]
    fills_sorted = sorted(trade_fills, key=lambda x: x.get("ts") or 0, reverse=True)
    symbol = ""
    fee_total = 0.0
    realized_pnl = 0.0
    last_fee = None
    last_fee_ccy = ""
    for f in fills_sorted:
        if not symbol:
            symbol = f.get("symbol") or ""
        fee_total += _to_float(f.get("fee")) or 0.0
        realized_pnl += _to_float(f.get("exec_pnl")) or 0.0
        if last_fee is None:
            last_fee = f.get("fee")
            last_fee_ccy = f.get("fee_ccy") or ""
    funding_events = [e for e in non_trade_events if "fund" in (e.get("exec_type") or "")]
    settlement_events = [e for e in non_trade_events if e not in funding_events]

    def _event_amount(evt: dict) -> float:
        fee = _to_float(evt.get("fee"))
        if fee not in (None, 0.0):
            return fee
        return _to_float(evt.get("exec_pnl")) or 0.0

    funding_fee_total = sum(_event_amount(e) for e in funding_events)
    return {
        "ts": time.time(),
        "now_ms": time.time() * 1000.0,
        "symbol": symbol,
        "session_pnl": realized_pnl,
        "realized_pnl": realized_pnl,
        "unrealized_pnl": None,
        "fee_total": fee_total,
        "funding_fee_total": funding_fee_total or 0.0,
        "last_fee": last_fee,
        "last_fee_ccy": last_fee_ccy,
        "last_fee_estimated": False,
        "fee_rate_maker": None,
        "fee_rate_taker": None,
        "fee_rate_source": "bybit_api",
        "inventory": None,
        "fill_rate": None,
        "orders": [],
        "fills": fills_sorted,
        "funding_events": funding_events,
        "settlement_events": settlement_events,
        "round_trips": [],
        "glft": {},
    }


def _resolve_data_root(cli_value: str | None) -> Path:
    if cli_value:
        return Path(cli_value).expanduser().resolve()
    env_value = os.getenv("TRADE_OPTIMIZER_DATA_ROOT", "").strip()
    if env_value:
        return Path(env_value).expanduser().resolve()
    default_candidate = ROOT.parent / "Resilient Maker"
    if default_candidate.exists():
        return default_candidate.resolve()
    return ROOT.resolve()


def _read_json(path: Path) -> dict | None:
    try:
        if not path.exists():
            return None
        with path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


class TradeOptimizerHandler(SimpleHTTPRequestHandler):
    data_root: Path = ROOT
    mode: str = "file"  # file | bybit_api | sync | realtime
    bybit_key: str = ""
    bybit_secret: str = ""
    bybit_url: str = "https://api.bybit.com"
    bybit_category: str = "linear"
    bybit_client: Optional[BybitAPIClient] = None
    ws_client: Optional[BybitWebSocketClient] = None
    ws_thread: Optional[threading.Thread] = None
    trading_key: str = ""
    trading_secret: str = ""
    trading_mode: str = ""
    trading_demo: bool = False
    trading_client: Optional[BybitAPIClient] = None
    trading_service: Optional[TradingService] = None

    def handle_one_request(self) -> None:
        try:
            super().handle_one_request()
        except (ConnectionResetError, BrokenPipeError):
            return

    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=str(DASHBOARD_DIR), **kwargs)

    def end_headers(self):
        self.send_header("Cache-Control", "no-store")
        super().end_headers()

    def _write_json(self, payload: dict, status: int = 200):
        """Write a JSON response with safe serialization fallbacks."""
        payload = self._make_json_serializable(payload)

        body: bytes
        try:
            import orjson
            body = orjson.dumps(
                payload,
                option=orjson.OPT_SERIALIZE_NUMPY | orjson.OPT_NON_STR_KEYS,
            )
        except Exception:
            def json_default(obj):
                try:
                    cls = type(obj).__name__
                    if "Timestamp" in cls or "Datetime" in cls or "datetime" in cls.lower() or "date" in cls.lower():
                        if hasattr(obj, "isoformat"):
                            return obj.isoformat()
                        if hasattr(obj, "strftime"):
                            return obj.strftime("%Y-%m-%dT%H:%M:%S.%f")
                        return str(obj)
                except Exception:
                    pass

                try:
                    import numpy as np
                    if isinstance(obj, (np.integer,)):
                        return int(obj)
                    if isinstance(obj, (np.floating,)):
                        return float(obj)
                except Exception:
                    pass

                return str(obj)

            try:
                body = json.dumps(payload, default=json_default).encode("utf-8")
            except Exception as e:
                import traceback
                print(f"⚠️  JSON serialization error: {e}")
                print(traceback.format_exc())
                payload = self._force_convert_all(payload)
                body = json.dumps(payload, default=str).encode("utf-8")
    
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        try:
            self.wfile.write(body)
        except (BrokenPipeError, ConnectionResetError):
            return

    def _force_convert_all(self, obj):
        """Recursively convert all objects (especially Timestamps) into JSON-safe types."""
        if isinstance(obj, dict):
            return {k: self._force_convert_all(v) for k, v in obj.items()}
        if isinstance(obj, (list, tuple)):
            return [self._force_convert_all(v) for v in obj]

        if isinstance(obj, (str, int, float, bool)) or obj is None:
            return obj

        try:
            cls = type(obj).__name__
            if "Timestamp" in cls or "Datetime" in cls or "datetime" in cls.lower() or "date" in cls.lower():
                if hasattr(obj, "isoformat"):
                    return obj.isoformat()
                if hasattr(obj, "strftime"):
                    return obj.strftime("%Y-%m-%dT%H:%M:%S.%f")
                return str(obj)
        except Exception:
            pass

        try:
            import numpy as np
            if isinstance(obj, (np.integer,)):
                return int(obj)
            if isinstance(obj, (np.floating,)):
                return float(obj)
        except Exception:
            pass

        return str(obj)
    
    def _make_json_serializable(self, obj):
        """Recursively convert non-JSON-serializable objects to serializable types"""
        import pandas as pd
        import numpy as np
        from datetime import date, datetime
        
        # First check class name for Timestamp (most reliable)
        class_name = type(obj).__name__
        
        # Check for Timestamp objects by class name first (catches all variants)
        if 'Timestamp' in class_name:
            try:
                if hasattr(obj, 'isoformat'):
                    return obj.isoformat()
                elif hasattr(obj, 'strftime'):
                    return obj.strftime('%Y-%m-%dT%H:%M:%S.%f')
                else:
                    return str(obj)
            except:
                return str(obj)
        
        # Then check isinstance for specific types
        try:
            if hasattr(pd, 'Timestamp') and isinstance(obj, pd.Timestamp):
                return obj.isoformat()
        except:
            pass
        
        if isinstance(obj, datetime):
            return obj.isoformat()
        elif isinstance(obj, date):
            return obj.isoformat()
        
        try:
            if hasattr(pd, 'DatetimeIndex') and isinstance(obj, pd.DatetimeIndex):
                return [str(x) for x in obj]
        except:
            pass
        
        if isinstance(obj, (np.integer, np.int64, np.int32, np.int16, np.int8)):
            return int(obj)
        elif isinstance(obj, (np.floating, np.float64, np.float32, np.float16)):
            return float(obj)
        elif isinstance(obj, np.ndarray):
            return [self._make_json_serializable(item) for item in obj.tolist()]
        
        try:
            if isinstance(obj, pd.DataFrame):
                # Convert DataFrame to dict, ensuring ALL Timestamps are converted
                try:
                    # Method 1: Convert datetime columns explicitly
                    df = obj.copy()
                    for col in df.columns:
                        col_dtype = str(df[col].dtype)
                        if 'datetime' in col_dtype.lower() or 'timestamp' in col_dtype.lower():
                            # Convert entire column to string, handling NaT
                            df[col] = df[col].apply(lambda x: x.isoformat() if hasattr(x, 'isoformat') and pd.notna(x) else str(x) if pd.notna(x) else None)
                    
                    df_dict = df.to_dict(orient="records")
                    # Recursively convert all values in the dict (catches any remaining Timestamps)
                    return [self._make_json_serializable(item) for item in df_dict]
                except Exception as e:
                    # Method 2: Convert each value individually (more robust)
                    try:
                        df_dict = obj.to_dict(orient="records")
                        result = []
                        for record in df_dict:
                            converted = {}
                            for k, v in record.items():
                                # Force convert each value
                                converted[k] = self._make_json_serializable(v)
                            result.append(converted)
                        return result
                    except:
                        # Method 3: Convert entire DataFrame to string representation
                        return str(obj)
        except:
            pass
        
        if isinstance(obj, dict):
            # Ensure keys are JSON-serializable (e.g. Timestamps as keys)
            converted = {}
            for k, v in obj.items():
                key = self._make_json_serializable(k) if not isinstance(k, (str, int, float, bool)) else k
                converted[str(key)] = self._make_json_serializable(v)
            return converted
        elif isinstance(obj, (list, tuple)):
            return [self._make_json_serializable(item) for item in obj]
        
        try:
            if hasattr(pd, 'isna') and pd.isna(obj):
                return None
        except:
            pass
        
        if isinstance(obj, (str, int, float, bool)) or obj is None:
            return obj
        
        # Final fallback: check for datetime-like objects
        if 'datetime' in class_name.lower() or 'date' in class_name.lower():
            try:
                if hasattr(obj, 'isoformat'):
                    return obj.isoformat()
                elif hasattr(obj, 'strftime'):
                    return obj.strftime('%Y-%m-%dT%H:%M:%S.%f')
            except:
                pass
        
        # Ultimate fallback: convert to string
        return str(obj)

    def _public_get(self, path: str, params: dict) -> dict | None:
        query = urllib.parse.urlencode(sorted(params.items()))
        url = f"{self.bybit_url}{path}?{query}"
        req = urllib.request.Request(url, method="GET")
        try:
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read().decode("utf-8"))
        except Exception:
            return None
        if str(data.get("retCode")) != "0":
            return None
        return data

    def _get_trading_client(self) -> Optional[BybitAPIClient]:
        if not self.trading_key or not self.trading_secret:
            return None
        if not self.trading_client:
            self.trading_client = BybitAPIClient(
                self.trading_key, self.trading_secret, self.bybit_url, self.bybit_category
            )
        return self.trading_client

    def _get_trading_service(self) -> Optional[TradingService]:
        client = self._get_trading_client()
        if not client:
            return None
        if not self.trading_service:
            self.trading_service = TradingService(client, BROKER_ID, self.bybit_category)
        return self.trading_service

    def _get_position_for_symbol(self, symbol: str) -> Optional[dict]:
        if not symbol:
            return None
        if self.mode == "realtime":
            with _realtime_lock:
                pos = _realtime_data["positions"].get(symbol)
                if pos:
                    return pos
        positions = self._fetch_positions().get("positions", [])
        for pos in positions:
            if (pos.get("symbol") or pos.get("Symbol")) == symbol:
                return pos
        return None

    def _admin_enabled(self) -> bool:
        try:
            return ADMIN_CONFIG_PATH.exists()
        except Exception:
            return False

    def _read_admin_config(self) -> dict:
        data = _read_json(ADMIN_CONFIG_PATH)
        return data or {}

    def _send_file(self, path: Path, content_type: str):
        if not path.exists():
            self.send_error(404, "File not found")
            return
        try:
            data = path.read_bytes()
        except Exception:
            self.send_error(500, "Failed to read file")
            return
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Content-Disposition", f'attachment; filename="{path.name}"')
        self.end_headers()
        try:
            self.wfile.write(data)
        except (BrokenPipeError, ConnectionResetError):
            return

    def _build_report_summary(self, analysis: dict) -> list[str]:
        basic = analysis.get("basic") or {}
        perf = analysis.get("performance") or {}
        risk = analysis.get("risk") or {}
        return [
            f"Total trades: {basic.get('total_trades', 0)}",
            f"Total volume: {basic.get('total_volume', 0):.2f}",
            f"Net PnL: {basic.get('net_pnl', 0):.4f}",
            f"Win rate: {perf.get('win_rate', 0) * 100:.2f}%",
            f"Max drawdown: {perf.get('max_drawdown', 0):.4f}",
            f"Volatility (annual): {risk.get('volatility_annualised', 0):.4f}",
        ]

    def _generate_pdf_report(self, analysis: dict, title: str) -> Path | None:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
        except Exception:
            return None
        REPORTS_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"report_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
        path = REPORTS_DIR / filename
        c = canvas.Canvas(str(path), pagesize=letter)
        width, height = letter
        y = height - 60
        c.setFont("Helvetica-Bold", 16)
        c.drawString(60, y, title)
        y -= 30
        c.setFont("Helvetica", 11)
        for line in self._build_report_summary(analysis):
            c.drawString(60, y, line)
            y -= 18
            if y < 80:
                c.showPage()
                y = height - 60
        c.showPage()
        c.save()
        return path

    def _generate_pptx_report(self, analysis: dict, title: str) -> Path | None:
        try:
            from pptx import Presentation
        except Exception:
            return None
        REPORTS_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"report_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
        path = REPORTS_DIR / filename
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Account trading report"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Summary"
        body = slide2.shapes.placeholders[1].text_frame
        for line in self._build_report_summary(analysis):
            p = body.add_paragraph()
            p.text = line
        prs.save(str(path))
        return path


    def do_GET(self):
        # Parse query parameters
        if "?" in self.path:
            path, query_str = self.path.split("?", 1)
            query_params = urllib.parse.parse_qs(query_str)
        else:
            path = self.path
            query_params = {}
        
        # Status endpoint
        if path == "/api/status":
            # API-backed modes (no metrics.json on disk)
            if self.mode in ("bybit_api", "sync", "realtime"):
                ready = bool(self.bybit_key and self.bybit_secret)
                ws_connected = _realtime_data["running"] if self.mode == "realtime" else None

                # In realtime mode we want the UI to show the product/source even before fills arrive.
                if self.mode == "realtime":
                    has_metrics = True
                    product = "realtime_mode"
                else:
                    has_metrics = ready
                    product = "api_direct"

                payload = {
                    "ok": True,
                    "has_metrics": has_metrics,
                    "product": product,
                    "exchange": "BYBIT",
                    "mode": self.mode,
                    "category": self.bybit_category,
                    "ws_connected": ws_connected,
                }
                return self._write_json(payload, status=200)

            metrics_path = self.data_root / "runtime" / "metrics.json"
            metrics = _read_json(metrics_path) or {}
            product = metrics.get("product") or "unknown"
            exchange = metrics.get("exchange") or metrics.get("primary_exchange") or ""
            has_metrics = bool(metrics)
            payload = {
                "ok": True,
                "has_metrics": has_metrics,
                "product": product,
                "exchange": exchange,
                "data_root": str(self.data_root),
                "metrics_mtime": metrics_path.stat().st_mtime if metrics_path.exists() else None,
                "mode": "file",
            }
            return self._write_json(payload, status=200)

        # Admin status endpoint
        if path == "/api/admin/status":
            return self._write_json({"ok": self._admin_enabled()}, status=200)

        # Metrics endpoint
        if path == "/api/metrics":
            if self.mode == "realtime":
                metrics = self._realtime_metrics()
                if metrics is None:
                    return self._write_json({"ok": False, "error": "realtime mode not ready"}, status=502)
                return self._write_json(metrics, status=200)
            
            elif self.mode == "sync":
                metrics = self._sync_metrics(query_params)
                if metrics is None:
                    return self._write_json({"ok": False, "error": "sync mode not ready"}, status=502)
                return self._write_json(metrics, status=200)
            
            elif self.mode == "bybit_api":
                metrics = self._bybit_metrics()
                if metrics is None:
                    return self._write_json({"ok": False, "error": "bybit api not ready"}, status=502)
                return self._write_json(metrics, status=200)

            metrics_path = self.data_root / "runtime" / "metrics.json"
            metrics = _read_json(metrics_path)
            if metrics is None:
                return self._write_json({"ok": False, "error": "metrics not found"}, status=404)
            return self._write_json(metrics, status=200)

        # Analysis endpoint (for sync and realtime modes)
        if path == "/api/analysis/fills":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "analysis not available in this mode"}, status=400)
            try:
                fills, error = self._get_analysis_fills(query_params)
                if error:
                    return self._write_json({"ok": False, "error": error, "fills": []}, status=502)
                return self._write_json({"ok": True, "fills": fills}, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/analysis/fills: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Analysis export error: {error_msg}"}, status=500)

        if path == "/api/admin/report":
            if not self._admin_enabled():
                return self._write_json({"ok": False, "error": "admin disabled"}, status=403)
            fmt = (query_params.get("format") or ["pdf"])[0].lower()
            symbol = (query_params.get("symbol") or [""])[0]
            start = (query_params.get("start") or [""])[0]
            end = (query_params.get("end") or [""])[0]
            days = (query_params.get("days") or ["30"])[0]
            qp = {"days": [days], "symbol": [symbol]}
            if start:
                qp["start"] = [start]
            if end:
                qp["end"] = [end]
            try:
                analysis = self._get_analysis(qp)
            except Exception as e:
                return self._write_json({"ok": False, "error": f"analysis error: {e}"}, status=500)
            if not analysis:
                return self._write_json({"ok": False, "error": "analysis unavailable"}, status=502)
            title = f"ResTM Report {datetime.utcnow().strftime('%Y-%m-%d')}"
            if fmt == "pptx":
                path = self._generate_pptx_report(analysis, title)
                if not path:
                    return self._write_json({"ok": False, "error": "pptx generator not available"}, status=500)
                return self._send_file(path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            path = self._generate_pdf_report(analysis, title)
            if not path:
                return self._write_json({"ok": False, "error": "pdf generator not available"}, status=500)
            return self._send_file(path, "application/pdf")

        if path == "/api/orders/history":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "order history not available in this mode"}, status=400)
            try:
                if not self.bybit_key or not self.bybit_secret:
                    return self._write_json({"ok": False, "error": "API credentials not set"}, status=400)
                if not self.bybit_client:
                    self.bybit_client = BybitAPIClient(
                        self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                    )
                symbol_list = query_params.get("symbol", [None])
                symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)
                start_ms, end_ms = self._parse_date_range(query_params)
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30
                limit_list = query_params.get("limit", [None])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else None)
                if limit_val not in (None, ""):
                    limit = int(limit_val)
                else:
                    limit = 20000
                if end_ms is None:
                    end_ms = int(time.time() * 1000)
                if start_ms is None:
                    start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                orders = _fetch_order_history(limit, start_ms, end_ms, symbol)
                return self._write_json({"ok": True, "orders": orders}, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/orders/history: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Order history error: {error_msg}"}, status=500)

        if path == "/api/analysis/compare":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "analysis not available in this mode"}, status=400)
            try:
                if not self.bybit_key or not self.bybit_secret:
                    return self._write_json({"ok": False, "error": "API credentials not set"}, status=400)
                if not self.bybit_client:
                    self.bybit_client = BybitAPIClient(
                        self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                    )
                symbol_list = query_params.get("symbol", [None])
                symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)
                start_ms, end_ms = self._parse_date_range(query_params)
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30
                limit_list = query_params.get("limit", [None])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else None)
                if limit_val not in (None, ""):
                    limit = int(limit_val)
                else:
                    limit = 20000
                if end_ms is None:
                    end_ms = int(time.time() * 1000)
                if start_ms is None:
                    start_ms = end_ms - (days * 24 * 60 * 60 * 1000)

                t0 = time.time()
                fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)
                t1 = time.time()
                orders = _fetch_order_history(limit, start_ms, end_ms, symbol)
                t2 = time.time()

                def _non_empty_fields(rows: list[dict], keys: list[str]) -> int:
                    count = 0
                    for k in keys:
                        if any((r.get(k) not in (None, "", 0, 0.0)) for r in rows):
                            count += 1
                    return count

                fill_symbols = sorted({(f.get("symbol") or "") for f in fills if f.get("symbol")})
                order_symbols = sorted({(o.get("symbol") or "") for o in orders if o.get("symbol")})
                filled_orders = [o for o in orders if float(o.get("cumExecQty") or 0) > 0]

                payload = {
                    "ok": True,
                    "fills": {
                        "rows": len(fills),
                        "symbols": fill_symbols,
                        "fields_present": _non_empty_fields(fills, ["exec_pnl", "fee", "is_maker", "price", "qty", "symbol", "side"]),
                        "ms": int((t1 - t0) * 1000),
                    },
                    "orders": {
                        "rows": len(orders),
                        "filled_rows": len(filled_orders),
                        "symbols": order_symbols,
                        "fields_present": _non_empty_fields(orders, ["orderType", "orderStatus", "timeInForce", "avgPrice", "price", "cumExecQty", "symbol", "side"]),
                        "ms": int((t2 - t1) * 1000),
                    },
                }
                return self._write_json(payload, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/analysis/compare: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Compare error: {error_msg}"}, status=500)

        if path == "/api/analysis":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "analysis not available in this mode"}, status=400)
            
            try:
                analysis = self._get_analysis(query_params)
                if analysis is None:
                    return self._write_json({"ok": False, "error": "analysis not available"}, status=502)
                return self._write_json(analysis, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/analysis: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Analysis error: {error_msg}"}, status=500)

        # Positions endpoint
        if path == "/api/positions":
            if self.mode == "realtime":
                # In realtime mode, try WebSocket data first, then fallback to REST API
                realtime_positions = self._get_realtime_positions()
                positions_list = realtime_positions.get("positions", [])
                # If WebSocket has positions, use them; otherwise fetch via REST API
                if positions_list and len(positions_list) > 0:
                    return self._write_json(realtime_positions, status=200)
                # Fallback to REST API if WebSocket hasn't received positions yet
                positions = self._fetch_positions()
                # Also update realtime_data with fetched positions for future use
                if positions.get("positions"):
                    with _realtime_lock:
                        for pos in positions["positions"]:
                            symbol = pos.get("symbol") or pos.get("Symbol", "")
                            if symbol:
                                _realtime_data["positions"][symbol] = pos
                return self._write_json(positions, status=200)
            elif self.mode in ("sync", "bybit_api"):
                positions = self._fetch_positions()
                return self._write_json(positions, status=200)
            if self.mode == "file":
                metrics_path = self.data_root / "runtime" / "metrics.json"
                metrics = _read_json(metrics_path) or {}
                positions = metrics.get("positions") or metrics.get("open_positions") or metrics.get("active_positions") or []
                return self._write_json({"positions": positions}, status=200)
            return self._write_json({"ok": False, "error": "positions not available in this mode"}, status=400)

        if path == "/api/klines":
            symbol = query_params.get("symbol", [""])[0]
            interval = query_params.get("interval", ["1"])[0]
            limit = int(query_params.get("limit", ["300"])[0])
            if not symbol:
                return self._write_json({"ok": False, "error": "symbol is required"}, status=400)
            params = {
                "category": self.bybit_category,
                "symbol": symbol,
                "interval": interval,
                "limit": str(min(max(limit, 1), 1000)),
            }
            data = self._public_get("/v5/market/kline", params)
            if not data:
                return self._write_json({"ok": False, "error": "kline fetch failed"}, status=502)
            rows = (data.get("result") or {}).get("list") or []
            klines = []
            for row in reversed(rows):
                try:
                    ts_ms = int(row[0])
                    klines.append({
                        "time": int(ts_ms / 1000),
                        "open": float(row[1]),
                        "high": float(row[2]),
                        "low": float(row[3]),
                        "close": float(row[4]),
                        "volume": float(row[5]),
                    })
                except Exception:
                    continue
            return self._write_json({"ok": True, "klines": klines}, status=200)

        if path == "/api/trades/stream":
            symbol = query_params.get("symbol", [""])[0]
            if not symbol:
                self.send_response(400)
                self.end_headers()
                return
            self.send_response(200)
            self.send_header("Content-Type", "text/event-stream")
            self.send_header("Cache-Control", "no-cache")
            self.send_header("Connection", "keep-alive")
            self.end_headers()
            seen_keys = deque()
            seen_limit = 2000
            seen_set = set()

            def _fill_key(f: dict) -> str:
                return f"{f.get('order_id','')}|{f.get('ts','')}|{f.get('side','')}|{f.get('price','')}|{f.get('qty','')}"

            def _public_trade_key(t: dict) -> str:
                return f"{t.get('execId','')}|{t.get('time','')}|{t.get('side','')}|{t.get('price','')}|{t.get('size','')}"

            def _normalize_public_trade(t: dict) -> Optional[dict]:
                try:
                    ts_ms = int(t.get("time") or 0)
                    price = float(t.get("price") or 0)
                    qty = float(t.get("size") or 0)
                except (TypeError, ValueError):
                    return None
                if not ts_ms or not price:
                    return None
                return {
                    "ts": ts_ms / 1000,
                    "price": price,
                    "qty": qty,
                    "side": t.get("side") or "",
                    "symbol": symbol,
                }

            with _realtime_lock:
                fills_snapshot = list(_realtime_data["fills"])
            trades = [f for f in fills_snapshot if f.get("exec_type") == "trade" and f.get("symbol") == symbol]
            trades = trades[-1000:]
            if not trades:
                public = self._public_get("/v5/market/recent-trade", {"category": self.bybit_category, "symbol": symbol, "limit": "1000"})
                rows = (public or {}).get("result", {}).get("list", []) if public else []
                trades = []
                for t in rows:
                    norm = _normalize_public_trade(t)
                    if norm:
                        trades.append(norm)
            trades = sorted(trades, key=lambda x: x.get("ts") or 0)
            snapshot_payload = {"trades": trades}
            try:
                self.wfile.write(f"event: snapshot\ndata: {json.dumps(snapshot_payload)}\n\n".encode("utf-8"))
                self.wfile.flush()
                for f in trades:
                    key = _fill_key(f)
                    if key not in seen_set:
                        seen_set.add(key)
                        seen_keys.append(key)
            except (BrokenPipeError, ConnectionResetError):
                return

            try:
                while True:
                    time.sleep(0.5)
                    with _realtime_lock:
                        fills = list(_realtime_data["fills"])
                    for f in fills:
                        if f.get("exec_type") != "trade":
                            continue
                        if f.get("symbol") != symbol:
                            continue
                        key = _fill_key(f)
                        if key in seen_set:
                            continue
                        seen_set.add(key)
                        seen_keys.append(key)
                        if len(seen_keys) > seen_limit:
                            old = seen_keys.popleft()
                            if old in seen_set:
                                seen_set.remove(old)
                        payload = json.dumps({
                            "ts": f.get("ts"),
                            "price": f.get("price"),
                            "qty": f.get("qty"),
                            "side": f.get("side"),
                            "symbol": f.get("symbol"),
                        })
                        self.wfile.write(f"data: {payload}\n\n".encode("utf-8"))
                    if not fills:
                        public = self._public_get("/v5/market/recent-trade", {"category": self.bybit_category, "symbol": symbol, "limit": "50"})
                        rows = (public or {}).get("result", {}).get("list", []) if public else []
                        for t in reversed(rows):
                            key = _public_trade_key(t)
                            if key in seen_set:
                                continue
                            seen_set.add(key)
                            seen_keys.append(key)
                            if len(seen_keys) > seen_limit:
                                old = seen_keys.popleft()
                                if old in seen_set:
                                    seen_set.remove(old)
                            norm = _normalize_public_trade(t)
                            if not norm:
                                continue
                            self.wfile.write(f"data: {json.dumps(norm)}\n\n".encode("utf-8"))
                    self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                return

        if path == "/api/trading/status":
            return self._write_json(
                {
                    "enabled": bool(self.trading_key and self.trading_secret and self.trading_mode),
                    "mode": self.trading_mode or "",
                    "has_env": bool(self.bybit_key and self.bybit_secret),
                    "demo": bool(self.trading_demo),
                },
                status=200,
            )

        if path == "/api/skew":
            skew_path = self.data_root / "runtime" / "skew_metrics.json"
            metrics = _read_json(skew_path)
            if metrics is None:
                return self._write_json({"ok": False, "error": "skew metrics not found"}, status=404)
            return self._write_json(metrics, status=200)

        return super().do_GET()

    def do_POST(self):
        content_length = int(self.headers.get("Content-Length", "0"))
        raw = self.rfile.read(content_length) if content_length > 0 else b""
        try:
            payload = json.loads(raw.decode("utf-8")) if raw else {}
        except Exception:
            return self._write_json({"ok": False, "error": "Invalid JSON body"}, status=400)

        if self.path == "/api/trading/config":
            api_key = str(payload.get("api_key", "")).strip()
            api_secret = str(payload.get("api_secret", "")).strip()
            mode = str(payload.get("mode", "")).strip()
            use_env = bool(payload.get("use_env"))
            demo = bool(payload.get("demo"))
            # Remove any hidden whitespace characters from pasted keys.
            api_key = "".join(api_key.split())
            api_secret = "".join(api_secret.split())
            if not api_key or not api_secret:
                if use_env and self.bybit_key and self.bybit_secret:
                    api_key = self.bybit_key
                    api_secret = self.bybit_secret
                else:
                    return self._write_json({"ok": False, "error": "API key and secret are required"}, status=400)
            if mode not in ("close_only", "open_close"):
                return self._write_json({"ok": False, "error": "Invalid trading mode"}, status=400)

            cls = type(self)
            cls.trading_key = api_key
            cls.trading_secret = api_secret
            cls.trading_mode = mode
            cls.trading_demo = demo
            cls.trading_client = None
            cls.trading_service = None
            return self._write_json({"ok": True, "mode": mode}, status=200)

        if self.path == "/api/trading/disable":
            cls = type(self)
            cls.trading_key = ""
            cls.trading_secret = ""
            cls.trading_mode = ""
            cls.trading_demo = False
            cls.trading_client = None
            cls.trading_service = None
            return self._write_json({"ok": True}, status=200)

        if self.path == "/api/orders/list":
            symbol = str(payload.get("symbol", "")).strip() or None
            trading_client = self._get_trading_client()
            primary_client = self.bybit_client
            if not primary_client and self.bybit_key and self.bybit_secret:
                self.bybit_client = BybitAPIClient(
                    self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                )
                primary_client = self.bybit_client
            client = trading_client or primary_client
            if not client:
                return self._write_json({"ok": False, "error": "API client not configured"}, status=400)

            def _try_fetch(target_client: BybitAPIClient) -> tuple[list[dict], Optional[str]]:
                if symbol:
                    orders = target_client.fetch_active_orders(symbol=symbol)
                    return orders, target_client.last_error
                errors: list[str] = []
                if self.bybit_category == "linear":
                    orders = target_client.fetch_active_orders(settle_coin="USDT")
                    if orders:
                        return orders, None
                    if target_client.last_error:
                        errors.append(target_client.last_error)
                    orders = target_client.fetch_active_orders(settle_coin="USDC")
                    if orders:
                        return orders, None
                    if target_client.last_error:
                        errors.append(target_client.last_error)
                with _realtime_lock:
                    recent_symbols = []
                    for pos in _realtime_data["positions"].values():
                        sym = pos.get("symbol") or pos.get("Symbol")
                        if sym and sym not in recent_symbols:
                            recent_symbols.append(sym)
                    for f in reversed(list(_realtime_data["fills"])):
                        sym = f.get("symbol")
                        if sym and sym not in recent_symbols:
                            recent_symbols.append(sym)
                        if len(recent_symbols) >= 8:
                            break
                if not recent_symbols:
                    positions = self._fetch_positions().get("positions", [])
                    for pos in positions:
                        sym = pos.get("symbol") or pos.get("Symbol")
                        if sym and sym not in recent_symbols:
                            recent_symbols.append(sym)
                for sym in recent_symbols:
                    orders = target_client.fetch_active_orders(symbol=sym)
                    if orders:
                        return orders, None
                    if target_client.last_error:
                        errors.append(target_client.last_error)
                return [], errors[-1] if errors else None

            orders, err = _try_fetch(client)
            if not orders and err and primary_client and client is not primary_client:
                orders, err = _try_fetch(primary_client)
            if not orders and err:
                return self._write_json({"ok": False, "error": err}, status=502)
            return self._write_json({"ok": True, "orders": orders}, status=200)

        if self.path == "/api/orders/cancel":
            order_id = str(payload.get("order_id", "")).strip()
            symbol = str(payload.get("symbol", "")).strip()
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            if not order_id or not symbol:
                return self._write_json({"ok": False, "error": "order_id and symbol are required"}, status=400)
            client = self._get_trading_client()
            if not client:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            ok = client.cancel_order(order_id=order_id, symbol=symbol)
            if not ok and client.last_error:
                return self._write_json({"ok": False, "error": client.last_error}, status=502)
            return self._write_json({"ok": True}, status=200)

        if self.path == "/api/orders/create":
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            symbol = str(payload.get("symbol", "")).strip()
            side = str(payload.get("side", "")).strip()
            qty = float(payload.get("qty", 0) or 0)
            order_mode = str(payload.get("order_mode", "")).strip()
            price = payload.get("price")
            reduce_only = bool(payload.get("reduce_only"))
            if not symbol or side not in ("Buy", "Sell"):
                return self._write_json({"ok": False, "error": "symbol and side are required"}, status=400)
            if qty <= 0:
                return self._write_json({"ok": False, "error": "qty must be > 0"}, status=400)
            if self.trading_mode == "close_only" and not reduce_only:
                return self._write_json({"ok": False, "error": "close_only mode requires reduce_only"}, status=400)
            service = self._get_trading_service()
            if not service:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            if order_mode == "market":
                result = service.place_order(symbol, side, qty, "Market", reduce_only=reduce_only)
            elif order_mode == "best":
                result = service.place_best_bid_offer(symbol, side, qty, reduce_only=reduce_only)
            elif order_mode == "chase":
                result = service.chase_limit(symbol, side, qty, reduce_only=reduce_only)
            elif order_mode == "limit":
                if price is None:
                    return self._write_json({"ok": False, "error": "price required for limit"}, status=400)
                result = service.place_order(symbol, side, qty, "Limit", price=float(price), reduce_only=reduce_only)
            else:
                return self._write_json({"ok": False, "error": "Unsupported order mode"}, status=400)
            if not result.ok:
                return self._write_json({"ok": False, "error": result.error}, status=502)
            return self._write_json({"ok": True, "order_id": result.order_id}, status=200)

        if self.path == "/api/positions/close":
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            symbol = str(payload.get("symbol", "")).strip()
            if not symbol:
                return self._write_json({"ok": False, "error": "symbol is required"}, status=400)
            pos = self._get_position_for_symbol(symbol)
            if not pos:
                return self._write_json({"ok": False, "error": "Position not found"}, status=404)
            size = float(pos.get("size") or pos.get("Size") or 0)
            side = str(pos.get("side") or pos.get("Side") or "").strip()
            if size <= 0 or not side:
                return self._write_json({"ok": False, "error": "Position size invalid"}, status=400)
            service = self._get_trading_service()
            if not service:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            result = service.close_position_market(symbol, side, size)
            if not result.ok:
                return self._write_json({"ok": False, "error": result.error}, status=502)
            return self._write_json({"ok": True, "order_id": result.order_id}, status=200)

        if self.path == "/api/positions/tpsl":
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            symbol = str(payload.get("symbol", "")).strip()
            take_profit = payload.get("take_profit")
            stop_loss = payload.get("stop_loss")
            if not symbol:
                return self._write_json({"ok": False, "error": "symbol is required"}, status=400)
            if take_profit is None and stop_loss is None:
                return self._write_json({"ok": False, "error": "take_profit or stop_loss required"}, status=400)
            service = self._get_trading_service()
            if not service:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            tp_val = float(take_profit) if take_profit is not None else None
            sl_val = float(stop_loss) if stop_loss is not None else None
            result = service.set_tpsl(symbol, tp_val, sl_val)
            if not result.ok:
                return self._write_json({"ok": False, "error": result.error}, status=502)
            return self._write_json({"ok": True}, status=200)

        if self.path == "/api/positions/smart-tp":
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            symbol = str(payload.get("symbol", "")).strip()
            if not symbol:
                return self._write_json({"ok": False, "error": "symbol is required"}, status=400)
            pos = self._get_position_for_symbol(symbol)
            if not pos:
                return self._write_json({"ok": False, "error": "Position not found"}, status=404)
            size = float(pos.get("size") or pos.get("Size") or 0)
            side = str(pos.get("side") or pos.get("Side") or "").strip()
            if size <= 0 or not side:
                return self._write_json({"ok": False, "error": "Position size invalid"}, status=400)
            service = self._get_trading_service()
            if not service:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            result = service.smart_tp(symbol, side, size)
            if not result.ok:
                return self._write_json({"ok": False, "error": result.error}, status=502)
            return self._write_json({"ok": True, "order_id": result.order_id}, status=200)

        if self.path == "/api/positions/smart-sl":
            if not (self.trading_key and self.trading_secret and self.trading_mode):
                return self._write_json({"ok": False, "error": "Trading is not enabled"}, status=400)
            symbol = str(payload.get("symbol", "")).strip()
            if not symbol:
                return self._write_json({"ok": False, "error": "symbol is required"}, status=400)
            pos = self._get_position_for_symbol(symbol)
            if not pos:
                return self._write_json({"ok": False, "error": "Position not found"}, status=404)
            size = float(pos.get("size") or pos.get("Size") or 0)
            side = str(pos.get("side") or pos.get("Side") or "").strip()
            break_even = float(pos.get("breakEvenPrice") or pos.get("BreakEvenPrice") or 0)
            mark_price = float(pos.get("markPrice") or pos.get("MarkPrice") or 0)
            if size <= 0 or not side:
                return self._write_json({"ok": False, "error": "Position size invalid"}, status=400)
            if break_even <= 0 or mark_price <= 0:
                return self._write_json({"ok": False, "error": "Missing break-even or mark price"}, status=400)
            service = self._get_trading_service()
            if not service:
                return self._write_json({"ok": False, "error": "Trading client not configured"}, status=400)
            result = service.smart_sl(symbol, side, size, break_even, mark_price)
            if not result.ok:
                return self._write_json({"ok": False, "error": result.error}, status=502)
            return self._write_json({"ok": True}, status=200)

        return self._write_json({"ok": False, "error": "Unknown POST endpoint"}, status=404)

    def _bybit_metrics(self) -> dict | None:
        """Original bybit_api mode metrics"""
        if not self.bybit_key or not self.bybit_secret:
            return None
        fills = self._bybit_fetch_fills(limit=200)
        metrics = _build_metrics_from_fills(fills)
        metrics["product"] = "api_direct"
        metrics["exchange"] = "BYBIT"
        metrics["session_id"] = f"api_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        return metrics
    
    def _sync_metrics(self, query_params: dict) -> dict | None:
        """Sync mode: fetch historical data and build metrics"""
        if not self.bybit_key or not self.bybit_secret:
            return None
        
        if not self.bybit_client:
            self.bybit_client = BybitAPIClient(
                self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
            )
        
        # Get parameters
        limit = int(query_params.get("limit", ["1000"])[0])
        symbol = query_params.get("symbol", [None])[0]
        days = int(query_params.get("days", ["30"])[0])
        
        # Calculate time range
        end_time = int(time.time() * 1000)
        start_time = end_time - (days * 24 * 60 * 60 * 1000)
        
        # Fetch fills (may need multiple requests)
        all_fills = []
        current_end = end_time
        max_per_request = 200
        
        max_range_ms = 7 * 24 * 60 * 60 * 1000
        while len(all_fills) < limit and current_end > start_time:
            window_start = max(start_time, current_end - max_range_ms)
            batch = self.bybit_client.fetch_executions(
                limit=max_per_request,
                symbol=symbol,
                start_time=window_start,
                end_time=current_end
            )
            if not batch:
                break
            all_fills.extend(batch)
            if len(batch) < max_per_request:
                current_end = window_start - 1
                continue
            # Move end_time to oldest fill timestamp
            oldest_ts = int(min(f["ts"] * 1000 for f in batch))
            if oldest_ts >= current_end:
                break
            current_end = oldest_ts - 1
            if current_end <= start_time:
                break
        
        # Limit to requested amount
        all_fills = all_fills[:limit]

        if not all_fills and self.bybit_client and self.bybit_client.last_error:
            return {
                "ok": False,
                "error": self.bybit_client.last_error,
                "mode": "sync",
            }
        
        # Build metrics
        metrics = _build_metrics_from_fills(all_fills)
        metrics["product"] = "sync_mode"
        metrics["exchange"] = "BYBIT"
        metrics["mode"] = "sync"
        metrics["session_id"] = f"sync_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        metrics["data_range_days"] = days
        metrics["total_fills_fetched"] = len(all_fills)
        
        return metrics
    
    def _realtime_metrics(self) -> dict | None:
        """Realtime mode: get metrics from live WebSocket data"""
        with _realtime_lock:
            fills = list(_realtime_data["fills"])
            non_trade_events = list(_realtime_data["non_trade_events"])
            last_update = _realtime_data["last_update"]
        
        # Return REST-backed metrics if websocket is connected but no fills yet
        if not fills and self.bybit_key and self.bybit_secret:
            if not self.bybit_client:
                self.bybit_client = BybitAPIClient(
                    self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                )
            fallback_fills = self.bybit_client.fetch_executions(limit=200)
            if fallback_fills:
                metrics = _build_metrics_from_fills(fallback_fills)
                metrics["product"] = "realtime_mode"
                metrics["exchange"] = "BYBIT"
                metrics["mode"] = "realtime"
                metrics["session_id"] = f"realtime_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
                metrics["session_start_ts"] = SERVER_START_TS
                metrics["ws_connected"] = _realtime_data["running"]
                metrics["last_update"] = last_update
                metrics["message"] = "Backfilled from REST because no websocket fills yet."
                return metrics

        # Return empty metrics structure if no fills yet
        if not fills:
            return {
                "ts": time.time(),
                "now_ms": time.time() * 1000.0,
                "symbol": "",
                "session_pnl": 0.0,
                "realized_pnl": 0.0,
                "unrealized_pnl": None,
                "fee_total": 0.0,
                "funding_fee_total": 0.0,
                "last_fee": None,
                "last_fee_ccy": "",
                "last_fee_estimated": False,
                "fee_rate_maker": None,
                "fee_rate_taker": None,
                "fee_rate_source": "bybit_api",
                "inventory": None,
                "fill_rate": None,
                "orders": [],
                "fills": [],
                "funding_events": [],
                "settlement_events": [],
                "round_trips": [],
                "glft": {},
                "product": "realtime_mode",
                "exchange": "BYBIT",
                "mode": "realtime",
                "session_id": f"realtime_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}",
                "session_start_ts": SERVER_START_TS,
                "ws_connected": _realtime_data["running"],
                "last_update": last_update,
                "message": "Waiting for trade data from WebSocket...",
            }
        
        metrics = _build_metrics_from_fills(fills + non_trade_events)
        metrics["product"] = "realtime_mode"
        metrics["exchange"] = "BYBIT"
        metrics["mode"] = "realtime"
        metrics["session_id"] = f"realtime_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        metrics["ws_connected"] = _realtime_data["running"]
        metrics["last_update"] = last_update
        
        return metrics

    def _parse_date_range(self, query_params: dict) -> tuple[int | None, int | None]:
        start_list = query_params.get("start", [""])
        end_list = query_params.get("end", [""])
        start_str = start_list[0] if isinstance(start_list, list) and start_list else (start_list or "")
        end_str = end_list[0] if isinstance(end_list, list) and end_list else (end_list or "")
        if not start_str and not end_str:
            return None, None
        start_ms = None
        end_ms = None
        try:
            if start_str:
                start_dt = datetime.strptime(start_str, "%Y-%m-%d").replace(tzinfo=timezone.utc)
                start_ms = int(start_dt.timestamp() * 1000)
            if end_str:
                end_dt = datetime.strptime(end_str, "%Y-%m-%d").replace(tzinfo=timezone.utc)
                end_dt = end_dt + timedelta(days=1) - timedelta(milliseconds=1)
                end_ms = int(end_dt.timestamp() * 1000)
        except Exception:
            return None, None
        return start_ms, end_ms

    def _merge_fills_unique(self, base: list[dict], extra: list[dict]) -> list[dict]:
        def _key(f: dict) -> tuple:
            ts_val = _normalize_ts(f.get("ts") or f.get("timestamp"))
            ts_ms = int(ts_val * 1000) if ts_val else 0
            qty = round(float(f.get("qty") or 0), 10)
            price = round(float(f.get("price") or 0), 10)
            return (
                f.get("order_id") or "",
                f.get("side") or "",
                f.get("symbol") or "",
                ts_ms,
                qty,
                price,
            )

        seen = set()
        merged: list[dict] = []
        for f in base:
            k = _key(f)
            if k in seen:
                continue
            seen.add(k)
            merged.append(f)
        for f in extra:
            k = _key(f)
            if k in seen:
                continue
            seen.add(k)
            merged.append(f)
        return merged

    def _get_analysis_fills(self, query_params: dict) -> tuple[list[dict], str | None]:
        fills: list[dict] = []
        days = 30
        symbol_list = query_params.get("symbol", [None])
        symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)
        source_list = query_params.get("source", [""])
        source = source_list[0] if isinstance(source_list, list) and source_list else (source_list or "")
        force_rest = str(source).lower() == "rest"
        start_ms, end_ms = self._parse_date_range(query_params)
        has_range = start_ms is not None or end_ms is not None

        def _fetch_rest_fills(limit: int, start_time: int, end_time: int, symbol_filter: str | None) -> list[dict]:
            results: list[dict] = []
            page_limit = min(200, max(1, limit))
            max_range_ms = 7 * 24 * 60 * 60 * 1000
            window_start = start_time
            while len(results) < limit and window_start <= end_time:
                window_end = min(end_time, window_start + max_range_ms - 1)
                cursor: str | None = None
                window_count = 0
                while len(results) < limit:
                    batch, cursor = self.bybit_client.fetch_executions_page(
                        limit=page_limit,
                        symbol=symbol_filter,
                        start_time=window_start,
                        end_time=window_end,
                        cursor=cursor
                    )
                    if not batch:
                        break
                    results.extend(batch)
                    window_count += len(batch)
                    if not cursor:
                        break
                window_start = window_end + 1
            return results[:limit]

        def _fetch_closed_pnl(limit: int, start_time: int, end_time: int, symbol_filter: str | None) -> list[dict]:
            results: list[dict] = []
            page_limit = min(200, max(1, limit))
            max_range_ms = 7 * 24 * 60 * 60 * 1000
            window_start = start_time
            while len(results) < limit and window_start <= end_time:
                window_end = min(end_time, window_start + max_range_ms - 1)
                cursor: str | None = None
                while len(results) < limit:
                    batch, cursor = self.bybit_client.fetch_closed_pnl_page(
                        limit=page_limit,
                        symbol=symbol_filter,
                        start_time=window_start,
                        end_time=window_end,
                        cursor=cursor
                    )
                    if not batch:
                        break
                    results.extend(batch)
                    if not cursor:
                        break
                window_start = window_end + 1
            return results[:limit]

        def _fetch_order_history(limit: int, start_time: int, end_time: int, symbol_filter: str | None) -> list[dict]:
            results: list[dict] = []
            page_limit = min(200, max(1, limit))
            max_range_ms = 7 * 24 * 60 * 60 * 1000
            window_start = start_time
            while len(results) < limit and window_start <= end_time:
                window_end = min(end_time, window_start + max_range_ms - 1)
                cursor: str | None = None
                while len(results) < limit:
                    batch, cursor = self.bybit_client.fetch_order_history_page(
                        limit=page_limit,
                        symbol=symbol_filter,
                        start_time=window_start,
                        end_time=window_end,
                        cursor=cursor
                    )
                    if not batch:
                        break
                    results.extend(batch)
                    if not cursor:
                        break
                window_start = window_end + 1
            return results[:limit]

        try:
            if force_rest:
                if not self.bybit_key or not self.bybit_secret:
                    return [], "API credentials not set for REST export"
                if not self.bybit_client:
                    self.bybit_client = BybitAPIClient(
                        self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                    )
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30
                limit_list = query_params.get("limit", [None])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else None)
                if limit_val not in (None, ""):
                    limit = int(limit_val)
                else:
                    limit = 20000 if days >= 14 or has_range else 5000

                if end_ms is None:
                    end_ms = int(time.time() * 1000)
                if start_ms is None:
                    start_ms = end_ms - (days * 24 * 60 * 60 * 1000)

                fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)
                if symbol:
                    fills = [f for f in fills if f.get("symbol") == symbol]
                trade_fills = [f for f in fills if not f.get("exec_type") or f.get("exec_type") == "trade"]
                if trade_fills and all((f.get("exec_pnl") or 0) == 0 for f in trade_fills):
                    closed_rows = _fetch_closed_pnl(limit, start_ms, end_ms, symbol)
                    pnl_by_order: dict[str, float] = {}
                    qty_by_order: dict[str, float] = {}
                    for row in closed_rows:
                        oid = row.get("order_id") or ""
                        if not oid:
                            continue
                        pnl_by_order[oid] = pnl_by_order.get(oid, 0.0) + float(row.get("closed_pnl") or 0.0)
                        qty_by_order[oid] = qty_by_order.get(oid, 0.0) + float(row.get("qty") or 0.0)
                    if pnl_by_order:
                        for f in trade_fills:
                            oid = f.get("order_id") or ""
                            if oid not in pnl_by_order:
                                continue
                            total_qty = qty_by_order.get(oid) or 0.0
                            if total_qty <= 0:
                                continue
                            fill_qty = float(f.get("qty") or 0.0)
                            f["exec_pnl"] = pnl_by_order[oid] * (fill_qty / total_qty)
                return trade_fills, None

            if self.mode == "realtime":
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30
                limit_list = query_params.get("limit", [None])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else None)
                if limit_val not in (None, ""):
                    limit = int(limit_val)
                else:
                    limit = 20000 if has_range else 1000

                if force_rest:
                    if not self.bybit_client:
                        self.bybit_client = BybitAPIClient(
                            self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                        )
                    if end_ms is None:
                        end_ms = int(time.time() * 1000)
                    if start_ms is None:
                        start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                    fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)
                else:
                    with _realtime_lock:
                        fills = list(_realtime_data["fills"])

                if not force_rest and has_range:
                    if end_ms is None:
                        end_ms = int(time.time() * 1000)
                    if start_ms is None:
                        start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                    if not self.bybit_client:
                        self.bybit_client = BybitAPIClient(
                            self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                        )
                    rest_fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)
                    start_sec = start_ms / 1000 if start_ms else None
                    end_sec = end_ms / 1000 if end_ms else None
                    realtime_scoped = [f for f in fills if (start_sec is None or _normalize_ts(f.get("ts")) >= start_sec)
                                       and (end_sec is None or _normalize_ts(f.get("ts")) <= end_sec)]
                    fills = self._merge_fills_unique(rest_fills, realtime_scoped)
                elif not force_rest and days > 0:
                    cutoff = time.time() - (days * 24 * 60 * 60)
                    fills = [f for f in fills if _normalize_ts(f.get("ts")) >= cutoff]
                if symbol:
                    fills = [f for f in fills if f.get("symbol") == symbol]

                if not fills:
                    if not self.bybit_client:
                        self.bybit_client = BybitAPIClient(
                            self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                        )
                    if end_ms is None:
                        end_ms = int(time.time() * 1000)
                    if start_ms is None:
                        start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                    fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)
                    if not fills and self.bybit_client and self.bybit_client.last_error:
                        return [], self.bybit_client.last_error

            elif self.mode == "sync":
                if not self.bybit_client:
                    self.bybit_client = BybitAPIClient(
                        self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                    )
                limit_list = query_params.get("limit", [None])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else None)
                if limit_val not in (None, ""):
                    limit = int(limit_val)
                else:
                    limit = 20000 if has_range else 1000

                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30

                symbol_list = query_params.get("symbol", [None])
                symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)

                if end_ms is None:
                    end_ms = int(time.time() * 1000)
                if start_ms is None:
                    start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                fills = _fetch_rest_fills(limit, start_ms, end_ms, symbol)

            elif self.mode == "file":
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30

                metrics_path = self.data_root / "runtime" / "metrics.json"
                metrics = _read_json(metrics_path) or {}
                fills = metrics.get("fills") or []

                if has_range:
                    if end_ms is None:
                        end_ms = int(time.time() * 1000)
                    if start_ms is None:
                        start_ms = end_ms - (days * 24 * 60 * 60 * 1000)
                    start_sec = start_ms / 1000 if start_ms else None
                    end_sec = end_ms / 1000 if end_ms else None
                    fills = [f for f in fills if (start_sec is None or _normalize_ts(f.get("ts") or f.get("timestamp")) >= start_sec)
                             and (end_sec is None or _normalize_ts(f.get("ts") or f.get("timestamp")) <= end_sec)]
                elif days > 0:
                    cutoff = time.time() - (days * 24 * 60 * 60)
                    fills = [f for f in fills if _normalize_ts(f.get("ts") or f.get("timestamp")) >= cutoff]
                if symbol:
                    fills = [f for f in fills if f.get("symbol") == symbol]
            else:
                fills = []

        except Exception as e:
            import traceback
            print(f"⚠️  Error fetching fills for analysis: {e}")
            print(traceback.format_exc())
            return [], f"Error fetching data: {str(e)}"

        trade_fills = [f for f in fills if not f.get("exec_type") or f.get("exec_type") == "trade"]
        return trade_fills, None

    def _get_analysis(self, query_params: dict) -> dict | None:
        """Get comprehensive analysis for sync/realtime/file modes"""
        trade_fills, error = self._get_analysis_fills(query_params)
        if error:
            return {
                "basic": {"total_trades": 0, "message": error},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "API error while fetching trades"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }

        if not trade_fills:
            return {
                "basic": {"total_trades": 0, "message": "No trade data available yet. Waiting for trades..."},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "No trades to analyze"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }

        try:
            normalized_fills = []
            for f in trade_fills:
                if not isinstance(f, dict):
                    continue
                ts = _normalize_ts(f.get("ts") or f.get("timestamp"))
                copy = dict(f)
                copy["ts"] = ts
                normalized_fills.append(copy)
            analyzer = TradeAnalyzer(normalized_fills)
            analysis = analyzer.get_comprehensive_analysis()
            analysis = self._serialize_analysis(analysis)
            return analysis
        except Exception as e:
            import traceback
            print(f"⚠️  Error analyzing trades: {e}")
            print(traceback.format_exc())
            return {
                "basic": {"total_trades": len(trade_fills), "message": f"Analysis error: {str(e)}"},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "Error during analysis"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }
    
    def _serialize_analysis(self, analysis: dict) -> dict:
        """Convert non-JSON-serializable objects (like pandas Timestamps, numpy types) to JSON-compatible types"""
        # Use the same serialization function as _make_json_serializable
        return self._make_json_serializable(analysis)
    
    def _get_realtime_positions(self) -> dict:
        """Get current positions from realtime data"""
        with _realtime_lock:
            # Convert dict to list for consistency with API response
            positions_dict = dict(_realtime_data["positions"])
            positions_list = list(positions_dict.values()) if positions_dict else []
            return {"positions": positions_list, "last_update": _realtime_data["last_update"]}
    
    def _fetch_positions(self) -> dict:
        """Fetch positions via API"""
        if not self.bybit_client:
            if not self.bybit_key or not self.bybit_secret:
                return {"positions": [], "error": "API credentials not set"}
            self.bybit_client = BybitAPIClient(
                self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
            )

        collected: list[dict] = []

        def _merge_positions(items: list[dict]):
            for item in items or []:
                symbol = item.get("symbol") or item.get("Symbol")
                if not symbol:
                    continue
                if not any((p.get("symbol") or p.get("Symbol")) == symbol for p in collected):
                    collected.append(item)

        _merge_positions(self.bybit_client.fetch_positions())
        if self.bybit_category == "linear":
            _merge_positions(self.bybit_client.fetch_positions(settle_coin="USDT"))
            _merge_positions(self.bybit_client.fetch_positions(settle_coin="USDC"))

        if not collected:
            with _realtime_lock:
                recent_symbols = []
                for f in reversed(list(_realtime_data["fills"])):
                    sym = f.get("symbol")
                    if sym and sym not in recent_symbols:
                        recent_symbols.append(sym)
                    if len(recent_symbols) >= 5:
                        break
            for sym in recent_symbols:
                _merge_positions(self.bybit_client.fetch_positions(symbol=sym))
                if self.bybit_category == "linear":
                    _merge_positions(self.bybit_client.fetch_positions(settle_coin="USDT", symbol=sym))
                    _merge_positions(self.bybit_client.fetch_positions(settle_coin="USDC", symbol=sym))

        if not collected and self.bybit_client.last_error:
            return {"positions": [], "error": self.bybit_client.last_error}
        return {"positions": collected}

    def _bybit_fetch_fills(self, limit: int = 200) -> list[dict]:
        params = {
            "category": self.bybit_category,
            "limit": str(limit),
        }
        data = self._bybit_get("/v5/execution/list", params)
        if not data:
            return []
        result = data.get("result") or {}
        rows = result.get("list") or []
        fills = []
        for r in rows:
            exec_type = (r.get("execType") or "").lower()
            ts_ms = _safe_int(r.get("execTime"))
            qty = _to_float(r.get("execQty")) or 0.0
            price = _to_float(r.get("execPrice")) or 0.0
            notional = _to_float(r.get("execValue")) or (qty * price if qty and price else 0.0)
            fills.append(
                {
                    "exec_type": exec_type,
                    "ts": ts_ms / 1000 if ts_ms else 0,
                    "order_id": r.get("orderId") or "",
                    "side": r.get("side") or "",
                    "qty": qty,
                    "price": price,
                    "notional": notional,
                    "fee": _to_float(r.get("execFee")) or 0.0,
                    "fee_ccy": r.get("feeCurrency") or "",
                    "is_tp": False,
                    "symbol": r.get("symbol") or "",
                    "exec_pnl": _to_float(r.get("execPnl")) or 0.0,
                    "is_maker": bool(r.get("isMaker")) if r.get("isMaker") is not None else None,
                }
            )
        return fills

    def _bybit_get(self, path: str, params: dict) -> dict | None:
        query = urllib.parse.urlencode(sorted(params.items()))
        timestamp = str(int(time.time() * 1000))
        recv_window = "5000"
        payload = f"{timestamp}{self.bybit_key}{recv_window}{query}"
        signature = hmac.new(
            self.bybit_secret.encode("utf-8"),
            payload.encode("utf-8"),
            hashlib.sha256
        ).hexdigest()

        url = f"{self.bybit_url}{path}?{query}"
        req = urllib.request.Request(url, method="GET")
        req.add_header("X-BAPI-API-KEY", self.bybit_key)
        req.add_header("X-BAPI-SIGN", signature)
        req.add_header("X-BAPI-TIMESTAMP", timestamp)
        req.add_header("X-BAPI-RECV-WINDOW", recv_window)
        try:
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read().decode("utf-8"))
        except Exception:
            return None
        if str(data.get("retCode")) != "0":
            return None
        return data


def _normalize_execution(exec_data: dict) -> dict:
    """Normalize execution data from WebSocket to our format"""
    exec_type = (exec_data.get("execType") or "").lower()
    ts_ms = int(exec_data.get("execTime", 0))
    qty = float(exec_data.get("execQty", 0))
    price = float(exec_data.get("execPrice", 0))
    notional = float(exec_data.get("execValue", 0)) or (qty * price if qty and price else 0)
    
    return {
        "exec_type": exec_type,
        "ts": ts_ms / 1000 if ts_ms else time.time(),
        "order_id": exec_data.get("orderId") or "",
        "side": exec_data.get("side") or "",
        "qty": qty,
        "price": price,
        "notional": notional,
        "fee": float(exec_data.get("execFee", 0)),
        "fee_ccy": exec_data.get("feeCurrency") or "",
        "symbol": exec_data.get("symbol") or "",
        "exec_pnl": float(exec_data.get("execPnl", 0)),
        "is_maker": bool(exec_data.get("isMaker")) if exec_data.get("isMaker") is not None else None,
    }


def _start_realtime_websocket():
    """Start WebSocket connection in background thread for realtime mode"""
    if not TradeOptimizerHandler.bybit_key or not TradeOptimizerHandler.bybit_secret:
        print("⚠️  Cannot start realtime mode: API credentials not set")
        return
    
    testnet = "testnet" in TradeOptimizerHandler.bybit_url.lower()
    
    ws_client = BybitWebSocketClient(
        TradeOptimizerHandler.bybit_key,
        TradeOptimizerHandler.bybit_secret,
        TradeOptimizerHandler.bybit_category,
        testnet=testnet
    )
    
    # Register callbacks
    def on_execution(exec_data: dict):
        fill = _normalize_execution(exec_data)
        with _realtime_lock:
            if fill.get("exec_type") and fill.get("exec_type") != "trade":
                _realtime_data["non_trade_events"].append(fill)
            else:
                _realtime_data["fills"].append(fill)
            _realtime_data["last_update"] = time.time()
    
    def on_position(pos_data: dict):
        symbol = pos_data.get("symbol", "")
        if symbol:
            with _realtime_lock:
                _realtime_data["positions"][symbol] = pos_data
                _realtime_data["last_update"] = time.time()
    
    def on_order(order_data: dict):
        with _realtime_lock:
            _realtime_data["orders"].append(order_data)
            _realtime_data["last_update"] = time.time()
    
    ws_client.on_execution(on_execution)
    ws_client.on_position(on_position)
    ws_client.on_order(on_order)
    
    TradeOptimizerHandler.ws_client = ws_client
    
    def run_ws():
        _realtime_data["running"] = True
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(ws_client.run())
        except Exception as e:
            print(f"⚠️  WebSocket error: {e}")
        finally:
            _realtime_data["running"] = False
            loop.close()
    
    ws_thread = threading.Thread(target=run_ws, daemon=True)
    ws_thread.start()
    TradeOptimizerHandler.ws_thread = ws_thread
    print("✅ Realtime WebSocket started in background thread")


def main():
    parser = argparse.ArgumentParser(description="Trade Optimizer dashboard server")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8010)
    parser.add_argument("--data-root", default=None, help="Root directory containing runtime/metrics.json")
    parser.add_argument("--mode", default=None, help="file | bybit_api | sync | realtime (default: env TRADE_OPTIMIZER_MODE)")
    args = parser.parse_args()

    data_root = _resolve_data_root(args.data_root)
    TradeOptimizerHandler.data_root = data_root
    TradeOptimizerHandler.mode = (args.mode or _env("TRADE_OPTIMIZER_MODE", "file")).lower()
    TradeOptimizerHandler.bybit_key = _env("BYBIT_API_KEY")
    TradeOptimizerHandler.bybit_secret = _env("BYBIT_API_SECRET")
    TradeOptimizerHandler.bybit_url = _env("BYBIT_API_URL", "https://api.bybit.com")
    TradeOptimizerHandler.bybit_category = _env("BYBIT_CATEGORY", "linear")

    if not DASHBOARD_DIR.exists():
        raise SystemExit(f"Dashboard directory missing: {DASHBOARD_DIR}")

    # Start WebSocket for realtime mode
    if TradeOptimizerHandler.mode == "realtime":
        if not TradeOptimizerHandler.bybit_key or not TradeOptimizerHandler.bybit_secret:
            print("⚠️  WARNING: Realtime mode requires BYBIT_API_KEY and BYBIT_API_SECRET")
            print("   Falling back to sync mode")
            TradeOptimizerHandler.mode = "sync"
        else:
            # Initialize API client for fetching initial positions
            TradeOptimizerHandler.bybit_client = BybitAPIClient(
                TradeOptimizerHandler.bybit_key,
                TradeOptimizerHandler.bybit_secret,
                TradeOptimizerHandler.bybit_url,
                TradeOptimizerHandler.bybit_category
            )
            # Fetch initial positions via REST API
            try:
                initial_positions = TradeOptimizerHandler.bybit_client.fetch_positions()
                with _realtime_lock:
                    for pos in initial_positions:
                        symbol = pos.get("symbol") or pos.get("Symbol", "")
                        if symbol:
                            _realtime_data["positions"][symbol] = pos
                if initial_positions:
                    print(f"✅ Loaded {len(initial_positions)} initial positions via REST API")
            except Exception as e:
                print(f"⚠️  Could not fetch initial positions: {e}")
            
            _start_realtime_websocket()
            time.sleep(2)  # Give WebSocket time to connect

    server = HTTPServer((args.host, args.port), TradeOptimizerHandler)
    print(f"Trade Optimizer running at http://{args.host}:{args.port}/")
    print(f"Data root: {data_root}")
    print(f"Mode: {TradeOptimizerHandler.mode}")
    if TradeOptimizerHandler.mode == "realtime":
        print("   Realtime mode: WebSocket connected for live data")
    elif TradeOptimizerHandler.mode == "sync":
        print("   Sync mode: Historical data analysis via REST API")
    
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n🛑 Shutting down...")
        if TradeOptimizerHandler.ws_client:
            asyncio.run(TradeOptimizerHandler.ws_client.close())


if __name__ == "__main__":
    main()

    def _make_json_serializable(self, obj):
        """Recursively convert non-JSON-serializable objects to serializable types"""
        import pandas as pd
        import numpy as np
        from datetime import date, datetime
        
        # First check class name for Timestamp (most reliable)
        class_name = type(obj).__name__
        
        # Check for Timestamp objects by class name first (catches all variants)
        if 'Timestamp' in class_name:
            try:
                if hasattr(obj, 'isoformat'):
                    return obj.isoformat()
                elif hasattr(obj, 'strftime'):
                    return obj.strftime('%Y-%m-%dT%H:%M:%S.%f')
                else:
                    return str(obj)
            except:
                return str(obj)
        
        # Then check isinstance for specific types
        try:
            if hasattr(pd, 'Timestamp') and isinstance(obj, pd.Timestamp):
                return obj.isoformat()
        except:
            pass
        
        if isinstance(obj, datetime):
            return obj.isoformat()
        elif isinstance(obj, date):
            return obj.isoformat()
        
        try:
            if hasattr(pd, 'DatetimeIndex') and isinstance(obj, pd.DatetimeIndex):
                return [str(x) for x in obj]
        except:
            pass
        
        if isinstance(obj, (np.integer, np.int64, np.int32, np.int16, np.int8)):
            return int(obj)
        elif isinstance(obj, (np.floating, np.float64, np.float32, np.float16)):
            return float(obj)
        elif isinstance(obj, np.ndarray):
            return [self._make_json_serializable(item) for item in obj.tolist()]
        
        try:
            if isinstance(obj, pd.DataFrame):
                # Convert DataFrame to dict, ensuring ALL Timestamps are converted
                try:
                    # Method 1: Convert datetime columns explicitly
                    df = obj.copy()
                    for col in df.columns:
                        col_dtype = str(df[col].dtype)
                        if 'datetime' in col_dtype.lower() or 'timestamp' in col_dtype.lower():
                            # Convert entire column to string, handling NaT
                            df[col] = df[col].apply(lambda x: x.isoformat() if hasattr(x, 'isoformat') and pd.notna(x) else str(x) if pd.notna(x) else None)
                    
                    df_dict = df.to_dict(orient="records")
                    # Recursively convert all values in the dict (catches any remaining Timestamps)
                    return [self._make_json_serializable(item) for item in df_dict]
                except Exception as e:
                    # Method 2: Convert each value individually (more robust)
                    try:
                        df_dict = obj.to_dict(orient="records")
                        result = []
                        for record in df_dict:
                            converted = {}
                            for k, v in record.items():
                                # Force convert each value
                                converted[k] = self._make_json_serializable(v)
                            result.append(converted)
                        return result
                    except:
                        # Method 3: Convert entire DataFrame to string representation
                        return str(obj)
        except:
            pass
        
        if isinstance(obj, dict):
            # Ensure keys are JSON-serializable (e.g. Timestamps as keys)
            converted = {}
            for k, v in obj.items():
                key = self._make_json_serializable(k) if not isinstance(k, (str, int, float, bool)) else k
                converted[str(key)] = self._make_json_serializable(v)
            return converted
        elif isinstance(obj, (list, tuple)):
            return [self._make_json_serializable(item) for item in obj]
        
        try:
            if hasattr(pd, 'isna') and pd.isna(obj):
                return None
        except:
            pass
        
        if isinstance(obj, (str, int, float, bool)) or obj is None:
            return obj
        
        # Final fallback: check for datetime-like objects
        if 'datetime' in class_name.lower() or 'date' in class_name.lower():
            try:
                if hasattr(obj, 'isoformat'):
                    return obj.isoformat()
                elif hasattr(obj, 'strftime'):
                    return obj.strftime('%Y-%m-%dT%H:%M:%S.%f')
            except:
                pass
        
        # Ultimate fallback: convert to string
        return str(obj)

    def do_GET(self):
        # Parse query parameters
        if "?" in self.path:
            path, query_str = self.path.split("?", 1)
            query_params = urllib.parse.parse_qs(query_str)
        else:
            path = self.path
            query_params = {}
        
        # Status endpoint
        if path == "/api/status":
            # API-backed modes (no metrics.json on disk)
            if self.mode in ("bybit_api", "sync", "realtime"):
                ready = bool(self.bybit_key and self.bybit_secret)
                ws_connected = _realtime_data["running"] if self.mode == "realtime" else None

                # In realtime mode we want the UI to show the product/source even before fills arrive.
                if self.mode == "realtime":
                    has_metrics = True
                    product = "realtime_mode"
                else:
                    has_metrics = ready
                    product = "api_direct"

                payload = {
                    "ok": True,
                    "has_metrics": has_metrics,
                    "product": product,
                    "exchange": "BYBIT",
                    "mode": self.mode,
                    "category": self.bybit_category,
                    "ws_connected": ws_connected,
                }
                return self._write_json(payload, status=200)

            metrics_path = self.data_root / "runtime" / "metrics.json"
            metrics = _read_json(metrics_path) or {}
            product = metrics.get("product") or "unknown"
            exchange = metrics.get("exchange") or metrics.get("primary_exchange") or ""
            has_metrics = bool(metrics)
            payload = {
                "ok": True,
                "has_metrics": has_metrics,
                "product": product,
                "exchange": exchange,
                "data_root": str(self.data_root),
                "metrics_mtime": metrics_path.stat().st_mtime if metrics_path.exists() else None,
                "mode": "file",
            }
            return self._write_json(payload, status=200)

        # Metrics endpoint
        if path == "/api/metrics":
            if self.mode == "realtime":
                metrics = self._realtime_metrics()
                if metrics is None:
                    return self._write_json({"ok": False, "error": "realtime mode not ready"}, status=502)
                return self._write_json(metrics, status=200)
            
            elif self.mode == "sync":
                metrics = self._sync_metrics(query_params)
                if metrics is None:
                    return self._write_json({"ok": False, "error": "sync mode not ready"}, status=502)
                return self._write_json(metrics, status=200)
            
            elif self.mode == "bybit_api":
                metrics = self._bybit_metrics()
                if metrics is None:
                    return self._write_json({"ok": False, "error": "bybit api not ready"}, status=502)
                return self._write_json(metrics, status=200)

            metrics_path = self.data_root / "runtime" / "metrics.json"
            metrics = _read_json(metrics_path)
            if metrics is None:
                return self._write_json({"ok": False, "error": "metrics not found"}, status=404)
            return self._write_json(metrics, status=200)

        # Analysis endpoint (for sync and realtime modes)
        if path == "/api/analysis/fills":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "analysis not available in this mode"}, status=400)
            try:
                fills, error = self._get_analysis_fills(query_params)
                if error:
                    return self._write_json({"ok": False, "error": error, "fills": []}, status=502)
                return self._write_json({"ok": True, "fills": fills}, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/analysis/fills: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Analysis export error: {error_msg}"}, status=500)

        if path == "/api/analysis":
            if self.mode not in ("sync", "realtime", "file"):
                return self._write_json({"ok": False, "error": "analysis not available in this mode"}, status=400)
            
            try:
                analysis = self._get_analysis(query_params)
                if analysis is None:
                    return self._write_json({"ok": False, "error": "analysis not available"}, status=502)
                return self._write_json(analysis, status=200)
            except Exception as e:
                import traceback
                error_msg = str(e)
                print(f"⚠️  Error in /api/analysis: {error_msg}")
                print(traceback.format_exc())
                return self._write_json({"ok": False, "error": f"Analysis error: {error_msg}"}, status=500)

        # Positions endpoint
        if path == "/api/positions":
            if self.mode == "realtime":
                # In realtime mode, try WebSocket data first, then fallback to REST API
                realtime_positions = self._get_realtime_positions()
                positions_list = realtime_positions.get("positions", [])
                # If WebSocket has positions, use them; otherwise fetch via REST API
                if positions_list and len(positions_list) > 0:
                    return self._write_json(realtime_positions, status=200)
                # Fallback to REST API if WebSocket hasn't received positions yet
                positions = self._fetch_positions()
                # Also update realtime_data with fetched positions for future use
                if positions.get("positions"):
                    with _realtime_lock:
                        for pos in positions["positions"]:
                            symbol = pos.get("symbol") or pos.get("Symbol", "")
                            if symbol:
                                _realtime_data["positions"][symbol] = pos
                return self._write_json(positions, status=200)
            elif self.mode in ("sync", "bybit_api"):
                positions = self._fetch_positions()
                return self._write_json(positions, status=200)
            return self._write_json({"ok": False, "error": "positions only available in API modes"}, status=400)

        if path == "/api/skew":
            skew_path = self.data_root / "runtime" / "skew_metrics.json"
            metrics = _read_json(skew_path)
            if metrics is None:
                return self._write_json({"ok": False, "error": "skew metrics not found"}, status=404)
            return self._write_json(metrics, status=200)

        return super().do_GET()

    def _bybit_metrics(self) -> dict | None:
        """Original bybit_api mode metrics"""
        if not self.bybit_key or not self.bybit_secret:
            return None
        fills = self._bybit_fetch_fills(limit=200)
        metrics = _build_metrics_from_fills(fills)
        metrics["product"] = "api_direct"
        metrics["exchange"] = "BYBIT"
        metrics["session_id"] = f"api_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        return metrics
    
    def _sync_metrics(self, query_params: dict) -> dict | None:
        """Sync mode: fetch historical data and build metrics"""
        if not self.bybit_key or not self.bybit_secret:
            return None
        
        if not self.bybit_client:
            self.bybit_client = BybitAPIClient(
                self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
            )
        
        # Get parameters
        limit = int(query_params.get("limit", ["1000"])[0])
        symbol = query_params.get("symbol", [None])[0]
        days = int(query_params.get("days", ["30"])[0])
        
        # Calculate time range
        end_time = int(time.time() * 1000)
        start_time = end_time - (days * 24 * 60 * 60 * 1000)
        
        # Fetch fills (may need multiple requests)
        all_fills = []
        current_end = end_time
        max_per_request = 200
        
        while len(all_fills) < limit and current_end > start_time:
            batch = self.bybit_client.fetch_executions(
                limit=max_per_request,
                symbol=symbol,
                end_time=current_end
            )
            if not batch:
                break
            all_fills.extend(batch)
            if len(batch) < max_per_request:
                break
            # Move end_time to oldest fill timestamp
            current_end = int(min(f["ts"] * 1000 for f in batch))
            if current_end <= start_time:
                break
        
        # Limit to requested amount
        all_fills = all_fills[:limit]
        
        # Build metrics
        metrics = _build_metrics_from_fills(all_fills)
        metrics["product"] = "sync_mode"
        metrics["exchange"] = "BYBIT"
        metrics["mode"] = "sync"
        metrics["session_id"] = f"sync_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        metrics["data_range_days"] = days
        metrics["total_fills_fetched"] = len(all_fills)
        
        return metrics
    
    def _realtime_metrics(self) -> dict | None:
        """Realtime mode: get metrics from live WebSocket data"""
        with _realtime_lock:
            fills = list(_realtime_data["fills"])
            last_update = _realtime_data["last_update"]
        
        # Return empty metrics structure if no fills yet
        if not fills:
            return {
                "ts": time.time(),
                "now_ms": time.time() * 1000.0,
                "symbol": "",
                "session_pnl": 0.0,
                "realized_pnl": 0.0,
                "unrealized_pnl": None,
                "fee_total": 0.0,
                "funding_fee_total": None,
                "last_fee": None,
                "last_fee_ccy": "",
                "last_fee_estimated": False,
                "fee_rate_maker": None,
                "fee_rate_taker": None,
                "fee_rate_source": "bybit_api",
                "inventory": None,
                "fill_rate": None,
                "orders": [],
                "fills": [],
                "funding_events": [],
                "round_trips": [],
                "glft": {},
                "product": "realtime_mode",
                "exchange": "BYBIT",
                "mode": "realtime",
                "session_id": f"realtime_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}",
                "session_start_ts": SERVER_START_TS,
                "ws_connected": _realtime_data["running"],
                "last_update": last_update,
                "message": "Waiting for trade data from WebSocket...",
            }
        
        metrics = _build_metrics_from_fills(fills)
        metrics["product"] = "realtime_mode"
        metrics["exchange"] = "BYBIT"
        metrics["mode"] = "realtime"
        metrics["session_id"] = f"realtime_{time.strftime('%Y%m%d_%H%M%S', time.gmtime(SERVER_START_TS))}"
        metrics["session_start_ts"] = SERVER_START_TS
        metrics["ws_connected"] = _realtime_data["running"]
        metrics["last_update"] = last_update
        
        return metrics
    
    def _get_analysis(self, query_params: dict) -> dict | None:
        """Get comprehensive analysis for sync/realtime modes"""
        fills = []
        # Defaults so realtime mode never references undefined variables
        days = 30
        symbol_list = query_params.get("symbol", [None])
        symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)
        
        try:
            if self.mode == "realtime":
                # Realtime mode: we already have fills in-memory; no need for days/start_time pagination
                with _realtime_lock:
                    fills = list(_realtime_data["fills"])
                if symbol:
                    fills = [f for f in fills if f.get("symbol") == symbol]

            elif self.mode == "sync":
                if not self.bybit_client:
                    self.bybit_client = BybitAPIClient(
                        self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
                    )
                # Safely parse query parameters with defaults
                limit_list = query_params.get("limit", ["1000"])
                limit_val = limit_list[0] if isinstance(limit_list, list) and len(limit_list) > 0 else (limit_list if isinstance(limit_list, str) else "1000")
                limit = int(limit_val) if limit_val else 1000
                
                days_list = query_params.get("days", ["30"])
                days_val = days_list[0] if isinstance(days_list, list) and len(days_list) > 0 else (days_list if isinstance(days_list, str) else "30")
                days = int(days_val) if days_val else 30
                
                symbol_list = query_params.get("symbol", [None])
                symbol = symbol_list[0] if isinstance(symbol_list, list) and len(symbol_list) > 0 else (symbol_list if isinstance(symbol_list, str) else None)

                end_time = int(time.time() * 1000)
                start_time = end_time - (days * 24 * 60 * 60 * 1000)

                current_end = end_time
                max_per_request = 200
                while len(fills) < limit and current_end > start_time:
                    batch = self.bybit_client.fetch_executions(
                        limit=max_per_request,
                        symbol=symbol,
                        end_time=current_end
                    )
                    if not batch:
                        break
                    fills.extend(batch)
                    if len(batch) < max_per_request:
                        break
                    current_end = int(min(f["ts"] * 1000 for f in batch))
                    if current_end <= start_time:
                        break
                fills = fills[:limit]

            else:
                # File/bybit_api modes don't support analysis
                fills = []
        
        except Exception as e:
            import traceback
            print(f"⚠️  Error fetching fills for analysis: {e}")
            print(traceback.format_exc())
            # Return empty analysis structure on error
            return {
                "basic": {"total_trades": 0, "message": f"Error fetching data: {str(e)}"},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "Error occurred"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }
        
        if not fills:
            # Return empty analysis structure instead of None
            return {
                "basic": {"total_trades": 0, "message": "No trade data available yet. Waiting for trades..."},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "No trades to analyze"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }
        
        try:
            analyzer = TradeAnalyzer(fills)
            analysis = analyzer.get_comprehensive_analysis()
            # Convert pandas Timestamps and other non-serializable objects to strings
            analysis = self._serialize_analysis(analysis)
            return analysis
        except Exception as e:
            import traceback
            print(f"⚠️  Error analyzing trades: {e}")
            print(traceback.format_exc())
            # Return partial analysis on error
            return {
                "basic": {"total_trades": len(fills), "message": f"Analysis error: {str(e)}"},
                "performance": {},
                "risk": {},
                "maker_taker": {},
                "consistency": {"score": 0, "message": "Error during analysis"},
                "hourly_stats": [],
                "daily_stats": [],
                "symbol_stats": [],
            }
    
    def _serialize_analysis(self, analysis: dict) -> dict:
        """Convert non-JSON-serializable objects (like pandas Timestamps, numpy types) to JSON-compatible types"""
        # Use the same serialization function as _make_json_serializable
        return self._make_json_serializable(analysis)
    
    def _get_realtime_positions(self) -> dict:
        """Get current positions from realtime data"""
        with _realtime_lock:
            # Convert dict to list for consistency with API response
            positions_dict = dict(_realtime_data["positions"])
            positions_list = list(positions_dict.values()) if positions_dict else []
            return {"positions": positions_list, "last_update": _realtime_data["last_update"]}
    
    def _fetch_positions(self) -> dict:
        """Fetch positions via API"""
        if not self.bybit_client:
            if not self.bybit_key or not self.bybit_secret:
                return {"positions": [], "error": "API credentials not set"}
            self.bybit_client = BybitAPIClient(
                self.bybit_key, self.bybit_secret, self.bybit_url, self.bybit_category
            )
        
        positions = self.bybit_client.fetch_positions()
        return {"positions": positions}

    def _bybit_fetch_fills(self, limit: int = 200) -> list[dict]:
        params = {
            "category": self.bybit_category,
            "limit": str(limit),
        }
        data = self._bybit_get("/v5/execution/list", params)
        if not data:
            return []
        result = data.get("result") or {}
        rows = result.get("list") or []
        fills = []
        for r in rows:
            ts_ms = _safe_int(r.get("execTime"))
            qty = _to_float(r.get("execQty")) or 0.0
            price = _to_float(r.get("execPrice")) or 0.0
            notional = _to_float(r.get("execValue")) or (qty * price if qty and price else 0.0)
            fills.append(
                {
                    "ts": ts_ms / 1000 if ts_ms else 0,
                    "order_id": r.get("orderId") or "",
                    "side": r.get("side") or "",
                    "qty": qty,
                    "price": price,
                    "notional": notional,
                    "fee": _to_float(r.get("execFee")) or 0.0,
                    "fee_ccy": r.get("feeCurrency") or "",
                    "is_tp": False,
                    "symbol": r.get("symbol") or "",
                    "exec_pnl": _to_float(r.get("execPnl")) or 0.0,
                    "is_maker": bool(r.get("isMaker")) if r.get("isMaker") is not None else None,
                }
            )
        return fills

    def _bybit_get(self, path: str, params: dict) -> dict | None:
        query = urllib.parse.urlencode(sorted(params.items()))
        timestamp = str(int(time.time() * 1000))
        recv_window = "5000"
        payload = f"{timestamp}{self.bybit_key}{recv_window}{query}"
        signature = hmac.new(
            self.bybit_secret.encode("utf-8"),
            payload.encode("utf-8"),
            hashlib.sha256
        ).hexdigest()

        url = f"{self.bybit_url}{path}?{query}"
        req = urllib.request.Request(url, method="GET")
        req.add_header("X-BAPI-API-KEY", self.bybit_key)
        req.add_header("X-BAPI-SIGN", signature)
        req.add_header("X-BAPI-TIMESTAMP", timestamp)
        req.add_header("X-BAPI-RECV-WINDOW", recv_window)
        try:
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read().decode("utf-8"))
        except Exception:
            return None
        if str(data.get("retCode")) != "0":
            return None
        return data


def _normalize_execution(exec_data: dict) -> dict:
    """Normalize execution data from WebSocket to our format"""
    ts_ms = int(exec_data.get("execTime", 0))
    qty = float(exec_data.get("execQty", 0))
    price = float(exec_data.get("execPrice", 0))
    notional = float(exec_data.get("execValue", 0)) or (qty * price if qty and price else 0)
    
    return {
        "ts": ts_ms / 1000 if ts_ms else time.time(),
        "order_id": exec_data.get("orderId") or "",
        "side": exec_data.get("side") or "",
        "qty": qty,
        "price": price,
        "notional": notional,
        "fee": float(exec_data.get("execFee", 0)),
        "fee_ccy": exec_data.get("feeCurrency") or "",
        "symbol": exec_data.get("symbol") or "",
        "exec_pnl": float(exec_data.get("execPnl", 0)),
        "is_maker": bool(exec_data.get("isMaker")) if exec_data.get("isMaker") is not None else None,
    }


def _start_realtime_websocket():
    """Start WebSocket connection in background thread for realtime mode"""
    if not TradeOptimizerHandler.bybit_key or not TradeOptimizerHandler.bybit_secret:
        print("⚠️  Cannot start realtime mode: API credentials not set")
        return
    
    testnet = "testnet" in TradeOptimizerHandler.bybit_url.lower()
    
    ws_client = BybitWebSocketClient(
        TradeOptimizerHandler.bybit_key,
        TradeOptimizerHandler.bybit_secret,
        TradeOptimizerHandler.bybit_category,
        testnet=testnet
    )
    
    # Register callbacks
    def on_execution(exec_data: dict):
        fill = _normalize_execution(exec_data)
        with _realtime_lock:
            _realtime_data["fills"].append(fill)
            _realtime_data["last_update"] = time.time()
    
    def on_position(pos_data: dict):
        symbol = pos_data.get("symbol", "")
        if symbol:
            with _realtime_lock:
                _realtime_data["positions"][symbol] = pos_data
                _realtime_data["last_update"] = time.time()
    
    def on_order(order_data: dict):
        with _realtime_lock:
            _realtime_data["orders"].append(order_data)
            _realtime_data["last_update"] = time.time()
    
    ws_client.on_execution(on_execution)
    ws_client.on_position(on_position)
    ws_client.on_order(on_order)
    
    TradeOptimizerHandler.ws_client = ws_client
    
    def run_ws():
        _realtime_data["running"] = True
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(ws_client.run())
        except Exception as e:
            print(f"⚠️  WebSocket error: {e}")
        finally:
            _realtime_data["running"] = False
            loop.close()
    
    ws_thread = threading.Thread(target=run_ws, daemon=True)
    ws_thread.start()
    TradeOptimizerHandler.ws_thread = ws_thread
    print("✅ Realtime WebSocket started in background thread")


def main():
    parser = argparse.ArgumentParser(description="Trade Optimizer dashboard server")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8010)
    parser.add_argument("--data-root", default=None, help="Root directory containing runtime/metrics.json")
    parser.add_argument("--mode", default=None, help="file | bybit_api | sync | realtime (default: env TRADE_OPTIMIZER_MODE)")
    args = parser.parse_args()

    data_root = _resolve_data_root(args.data_root)
    TradeOptimizerHandler.data_root = data_root
    TradeOptimizerHandler.mode = (args.mode or _env("TRADE_OPTIMIZER_MODE", "file")).lower()
    TradeOptimizerHandler.bybit_key = _env("BYBIT_API_KEY")
    TradeOptimizerHandler.bybit_secret = _env("BYBIT_API_SECRET")
    TradeOptimizerHandler.bybit_url = _env("BYBIT_API_URL", "https://api.bybit.com")
    TradeOptimizerHandler.bybit_category = _env("BYBIT_CATEGORY", "linear")

    if not DASHBOARD_DIR.exists():
        raise SystemExit(f"Dashboard directory missing: {DASHBOARD_DIR}")

    # Start WebSocket for realtime mode
    if TradeOptimizerHandler.mode == "realtime":
        if not TradeOptimizerHandler.bybit_key or not TradeOptimizerHandler.bybit_secret:
            print("⚠️  WARNING: Realtime mode requires BYBIT_API_KEY and BYBIT_API_SECRET")
            print("   Falling back to sync mode")
            TradeOptimizerHandler.mode = "sync"
        else:
            # Initialize API client for fetching initial positions
            TradeOptimizerHandler.bybit_client = BybitAPIClient(
                TradeOptimizerHandler.bybit_key,
                TradeOptimizerHandler.bybit_secret,
                TradeOptimizerHandler.bybit_url,
                TradeOptimizerHandler.bybit_category
            )
            # Fetch initial positions via REST API
            try:
                initial_positions = TradeOptimizerHandler.bybit_client.fetch_positions()
                with _realtime_lock:
                    for pos in initial_positions:
                        symbol = pos.get("symbol") or pos.get("Symbol", "")
                        if symbol:
                            _realtime_data["positions"][symbol] = pos
                if initial_positions:
                    print(f"✅ Loaded {len(initial_positions)} initial positions via REST API")
            except Exception as e:
                print(f"⚠️  Could not fetch initial positions: {e}")
            
            _start_realtime_websocket()
            time.sleep(2)  # Give WebSocket time to connect

    server = HTTPServer((args.host, args.port), TradeOptimizerHandler)
    print(f"Trade Optimizer running at http://{args.host}:{args.port}/")
    print(f"Data root: {data_root}")
    print(f"Mode: {TradeOptimizerHandler.mode}")
    if TradeOptimizerHandler.mode == "realtime":
        print("   Realtime mode: WebSocket connected for live data")
    elif TradeOptimizerHandler.mode == "sync":
        print("   Sync mode: Historical data analysis via REST API")
    
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n🛑 Shutting down...")
        if TradeOptimizerHandler.ws_client:
            asyncio.run(TradeOptimizerHandler.ws_client.close())


if __name__ == "__main__":
    main()
